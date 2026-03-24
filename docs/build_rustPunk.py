"""
Generic ALAN(v) Literature Review presentation builder.

Clones template slides (preserving backgrounds, colors, styling) then
replaces placeholder text with content from a slide_data list.

Template slides (0-based index -> type):
  0: title         BLANK           — {{Title}}, {{subtitle}}, AutoHelm, slide#
  1: agenda        CUSTOM_2_1      — Agenda, {{section 1..5}}, AutoHelm, slide#
  2: abstract      CUSTOM_3        — {{Description}}, Abstract, AutoHelm, slide#
  3: section       CUSTOM_12       — {{Section}}, {{Subsection A/B}}, AutoHelm, slide#
  4: subsection    CUSTOM_12(bg)   — {{Subsection}}, {{paragraph 1/2}}, ALAN(v), slide#
  5: paragraph     CUSTOM_7        — {{subsection A}} - {{Paragraph X}}, {{Description}}, AutoHelm, slide#
  6: sources       CUSTOM_11_1_1   — 4x {{source N}} + {{DescriptionN}}{{LinkN}}, AutoHelm, slide#
  7: cost_analysis BLANK_1_1_1...  — Cost Analysis, 3x {{paragraph N Title}} + {{paragraph N}}, AutoHelm, slide#

Usage:
  python build_rustPunk.py <slides.yaml> [output.pptx]
"""

import copy
import sys
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation

DOCS = Path(__file__).resolve().parent
ROOT = DOCS.parent
TEMPLATE = DOCS / "rustPunk_template.pptx"
OUTPUT = DOCS / "Literature_Review_Presentation.pptx"
BRAND = "ALAN"

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ─── Slide cloning ──────────────────────────────────────────────────


def clone_slide(prs, src_slide):
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    src_cSld = src_slide._element.find(f'{{{P}}}cSld')
    new_cSld = new_slide._element.find(f'{{{P}}}cSld')
    for child in list(new_cSld):
        new_cSld.remove(child)
    for child in src_cSld:
        new_cSld.append(copy.deepcopy(child))
    src_clr = src_slide._element.find(f'{{{P}}}clrMapOvr')
    if src_clr is not None:
        old_clr = new_slide._element.find(f'{{{P}}}clrMapOvr')
        if old_clr is not None:
            new_slide._element.replace(old_clr, copy.deepcopy(src_clr))
        else:
            new_slide._element.append(copy.deepcopy(src_clr))
    return new_slide


# ─── XML-level text helpers ─────────────────────────────────────────


def _sp_text(sp):
    return ''.join(t.text or '' for t in sp.findall(f'.//{{{A}}}t'))


def _sp_set_text(sp, new_text):
    txBody = sp.find(f'{{{P}}}txBody')
    if txBody is None:
        return
    first_rPr = txBody.find(f'.//{{{A}}}rPr')
    rPr_copy = copy.deepcopy(first_rPr) if first_rPr is not None else None
    first_pPr = txBody.find(f'.//{{{A}}}pPr')
    pPr_copy = copy.deepcopy(first_pPr) if first_pPr is not None else None
    for p in txBody.findall(f'{{{A}}}p'):
        txBody.remove(p)
    for line in new_text.split('\n'):
        p = etree.SubElement(txBody, f'{{{A}}}p')
        if pPr_copy is not None:
            p.append(copy.deepcopy(pPr_copy))
        r = etree.SubElement(p, f'{{{A}}}r')
        if rPr_copy is not None:
            r.append(copy.deepcopy(rPr_copy))
        t = etree.SubElement(r, f'{{{A}}}t')
        t.text = line


def _get_shapes(slide):
    cSld = slide._element.find(f'{{{P}}}cSld')
    spTree = cSld.find(f'{{{P}}}spTree')
    return [el for el in spTree if etree.QName(el.tag).localname == 'sp']


def _replace_brand(slide):
    for sp in _get_shapes(slide):
        if _sp_text(sp).strip() == 'AutoHelm':
            _sp_set_text(sp, BRAND)


def _apply_bg_from(slide, donor_slide):
    """Copy the background from donor_slide's layout onto slide at the slide level."""
    donor_layout = donor_slide.slide_layout
    layout_cSld = donor_layout._element.find(f'{{{P}}}cSld')
    layout_bg = layout_cSld.find(f'{{{P}}}bg') if layout_cSld is not None else None
    if layout_bg is None:
        return
    cSld = slide._element.find(f'{{{P}}}cSld')
    old_bg = cSld.find(f'{{{P}}}bg')
    if old_bg is not None:
        cSld.remove(old_bg)
    cSld.insert(0, copy.deepcopy(layout_bg))


def _force_text_color(slide, scheme_val='lt1'):
    """Set solidFill on every rPr in the slide to the given scheme color."""
    for sp in _get_shapes(slide):
        for rPr in sp.findall(f'.//{{{A}}}rPr'):
            old_fill = rPr.find(f'{{{A}}}solidFill')
            if old_fill is not None:
                rPr.remove(old_fill)
            sf = etree.SubElement(rPr, f'{{{A}}}solidFill')
            etree.SubElement(sf, f'{{{A}}}schemeClr').set('val', scheme_val)


# ─── Fill functions (one per slide type) ─────────────────────────────


def fill_title(slide, data):
    """data keys: title, subtitle (optional)"""
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        if txt == '{{Title}}':
            _sp_set_text(sp, data.get('title', ''))
        elif txt == '{{subtitle}}' or txt.startswith('State'):
            sub = data.get('subtitle', '')
            if sub:
                _sp_set_text(sp, sub)
    _replace_brand(slide)


def fill_agenda(slide, data):
    """data keys: sections (list of strings, any length)

    Keeps only as many agenda rows as there are sections and redistributes
    them evenly over the vertical area originally occupied by all 5 placeholders.
    """
    sections = [s for s in data.get('sections', []) if s]  # drop blanks
    n = len(sections)

    # Collect section-text shapes and number shapes keyed by their template index
    sec_shapes = {}   # {1: sp, 2: sp, ...}
    num_shapes = {}   # {1: sp, 2: sp, ...}
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        for i in range(1, 6):
            if txt.lower() == f'{{{{section {i}}}}}':
                sec_shapes[i] = sp
            elif txt == f'{i}.':
                num_shapes[i] = sp

    # Determine the original vertical span (top of first → bottom of last)
    ys = [int(sec_shapes[i].find(f'.//{{{A}}}off').get('y')) for i in sorted(sec_shapes)]
    h = int(sec_shapes[1].find(f'.//{{{A}}}ext').get('cy')) if sec_shapes else 0
    y_top = ys[0]
    y_bot = ys[-1] + h          # bottom edge of last item
    total = y_bot - y_top       # full distributable height

    cSld = slide._element.find(f'{{{P}}}cSld')
    spTree = cSld.find(f'{{{P}}}spTree')

    # Remove shapes for unused slots
    for i in range(1, 6):
        if i > n:
            if i in sec_shapes:
                spTree.remove(sec_shapes[i])
            if i in num_shapes:
                spTree.remove(num_shapes[i])

    # Redistribute kept shapes evenly
    for idx, label in enumerate(sections):
        i = idx + 1
        if n == 1:
            new_y = y_top + (total - h) // 2
        else:
            new_y = y_top + idx * (total - h) // (n - 1)
        if i in sec_shapes:
            _sp_set_text(sec_shapes[i], label)
            sec_shapes[i].find(f'.//{{{A}}}off').set('y', str(new_y))
        if i in num_shapes:
            num_shapes[i].find(f'.//{{{A}}}off').set('y', str(new_y))

    _replace_brand(slide)


def fill_abstract(slide, data):
    """data keys: description"""
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        if txt == '{{Description}}':
            _sp_set_text(sp, data.get('description', ''))
    _replace_brand(slide)


def fill_section(slide, data):
    """data keys: section, subsections (list of strings)"""
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        if txt.lower() == '{{section}}':
            _sp_set_text(sp, data.get('section', ''))
        elif '{{Subsection' in txt or '{{subsection' in txt:
            subs = data.get('subsections', [])
            _sp_set_text(sp, '\n'.join(subs))
    _replace_brand(slide)


def fill_subsection(slide, data):
    """data keys: section, subsections (list of strings)

    Uses the dedicated subsection template (slide 4) which has
    {{Subsection}} for the title and {{paragraph N}} for bullet items.
    """
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        if txt.lower() == '{{subsection}}':
            _sp_set_text(sp, data.get('section', ''))
        elif '{{paragraph' in txt.lower():
            subs = data.get('subsections', [])
            _sp_set_text(sp, '\n'.join(subs))
    _replace_brand(slide)


def fill_paragraph(slide, data):
    """data keys: heading, description (list of strings or single string)

    When description is a list with multiple items, each item gets its own
    text box distributed horizontally across the original Description area.
    """
    desc = data.get('description', '')
    if isinstance(desc, str):
        desc = [desc]

    desc_sp = None
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        if '{{subsection' in txt and '{{Paragraph' in txt:
            _sp_set_text(sp, data.get('heading', ''))
        elif txt == '{{Description}}':
            desc_sp = sp

    if desc_sp is not None:
        if len(desc) <= 1:
            _sp_set_text(desc_sp, desc[0] if desc else '')
        else:
            _split_description(slide, desc_sp, desc)

    _replace_brand(slide)


_COL_GAP = 91440  # ~0.1 inch gap between columns


def _split_description(slide, desc_sp, items):
    """Replace desc_sp with N horizontally-distributed copies, one per item."""
    xfrm = desc_sp.find(f'.//{{{A}}}xfrm')
    off = xfrm.find(f'{{{A}}}off')
    ext = xfrm.find(f'{{{A}}}ext')
    base_x = int(off.get('x'))
    base_y = int(off.get('y'))
    total_w = int(ext.get('cx'))
    total_h = int(ext.get('cy'))

    n = len(items)
    gap_total = _COL_GAP * (n - 1)
    col_w = (total_w - gap_total) // n

    cSld = slide._element.find(f'{{{P}}}cSld')
    spTree = cSld.find(f'{{{P}}}spTree')

    for i, text in enumerate(items):
        col_sp = copy.deepcopy(desc_sp)
        # Update position and width
        col_xfrm = col_sp.find(f'.//{{{A}}}xfrm')
        col_off = col_xfrm.find(f'{{{A}}}off')
        col_ext = col_xfrm.find(f'{{{A}}}ext')
        col_off.set('x', str(base_x + i * (col_w + _COL_GAP)))
        col_off.set('y', str(base_y))
        col_ext.set('cx', str(col_w))
        col_ext.set('cy', str(total_h))
        # Give each clone a unique id
        nvSpPr = col_sp.find(f'{{{P}}}nvSpPr')
        cNvPr = nvSpPr.find(f'{{{P}}}cNvPr')
        old_id = int(cNvPr.get('id', '0'))
        cNvPr.set('id', str(old_id + 1000 + i))
        cNvPr.set('name', f'DescCol{i}')
        _sp_set_text(col_sp, text)
        spTree.append(col_sp)

    # Remove the original placeholder
    spTree.remove(desc_sp)


def fill_sources(slide, data):
    """data keys: sources (list of dicts with 'name', 'description', 'link')"""
    sources = data.get('sources', [])
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        for i in range(1, 5):
            if txt == f'{{{{source {i}}}}}':
                src = sources[i - 1] if i - 1 < len(sources) else {}
                _sp_set_text(sp, src.get('name', ''))
            elif txt == f'{{{{Description{i}}}}}{{{{Link{i}}}}}':
                src = sources[i - 1] if i - 1 < len(sources) else {}
                desc = src.get('description', '')
                link = src.get('link', '')
                combined = f'{desc}\n{link}' if link else desc
                _sp_set_text(sp, combined)
    _replace_brand(slide)


def fill_cost_analysis(slide, data):
    """data keys: columns (list of 3 dicts with 'title' and 'body')"""
    cols = data.get('columns', [{}, {}, {}])
    for sp in _get_shapes(slide):
        txt = _sp_text(sp).strip()
        for i in range(1, 4):
            if txt == f'{{{{paragraph {i} Title}}}}' or txt == f'{{{{paragraph {i} title}}}}':
                col = cols[i - 1] if i - 1 < len(cols) else {}
                _sp_set_text(sp, col.get('title', ''))
            elif txt == f'{{{{paragraph {i}}}}}':
                col = cols[i - 1] if i - 1 < len(cols) else {}
                _sp_set_text(sp, col.get('body', ''))
    _replace_brand(slide)


# ─── Type dispatch ───────────────────────────────────────────────────

TEMPLATE_INDEX = {
    'title': 0,
    'agenda': 1,
    'abstract': 2,
    'section': 3,
    'subsection': 4,
    'paragraph': 5,
    'sources': 6,
    'cost_analysis': 7,
}

FILLERS = {
    'title': fill_title,
    'agenda': fill_agenda,
    'abstract': fill_abstract,
    'section': fill_section,
    'subsection': fill_subsection,
    'paragraph': fill_paragraph,
    'sources': fill_sources,
    'cost_analysis': fill_cost_analysis,
}


# ─── Main ────────────────────────────────────────────────────────────


def build(slide_data, output=OUTPUT):
    prs = Presentation(str(TEMPLATE))
    src = list(prs.slides)

    for entry in slide_data:
        stype = entry['type']
        s = clone_slide(prs, src[TEMPLATE_INDEX[stype]])
        FILLERS[stype](s, entry)

    # Remove original template slides
    id_list = prs.slides._sldIdLst
    originals = list(id_list)[:len(src)]
    for sid in originals:
        prs.part.drop_rel(sid.rId)
        id_list.remove(sid)

    prs.save(str(output))
    print(f"Saved {len(prs.slides)} slides -> {output}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <slides.yaml> [output.pptx]")
        sys.exit(1)
    yaml_path = Path(sys.argv[1])
    with open(yaml_path, encoding="utf-8") as f:
        slide_data = yaml.safe_load(f)
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else OUTPUT
    build(slide_data, output=out)
